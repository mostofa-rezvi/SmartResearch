import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_deck(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # ── Premium Cohesive Color Palette ──────────────────────────────────────
    C_DARK_BG       = RGBColor(15, 23, 42)      # #0F172A Deep Navy / Slate
    C_DARK_CARD     = RGBColor(30, 41, 59)     # #1E293B Card Background
    C_DARK_CARD_ALT = RGBColor(24, 33, 47)     # #18212F Slightly Darker Card
    C_DARK_BORDER   = RGBColor(51, 65, 85)     # #334155 Slate Border

    C_LIGHT_BG      = RGBColor(248, 250, 252)  # #F8FAFC Crisp Canvas
    C_LIGHT_CARD    = RGBColor(255, 255, 255)  # #FFFFFF Pure Card
    C_LIGHT_BORDER  = RGBColor(226, 232, 240)  # #E2E8F0 Soft Border
    C_LIGHT_CARD_ALT= RGBColor(241, 245, 249)  # #F1F5F9 Slate Tint

    # Semantic Accent Colors
    C_PRIMARY       = RGBColor(37, 99, 235)    # #2563EB Vibrant Brand Blue
    C_PRIMARY_LIGHT = RGBColor(239, 246, 255)  # #EFF6FF Blue Tint
    C_CYAN          = RGBColor(6, 182, 212)    # #06B6D4 Electric Cyan
    C_CYAN_LIGHT    = RGBColor(236, 254, 255)  # #ECFEFF Cyan Tint
    C_EMERALD       = RGBColor(16, 185, 129)   # #10B981 Vibrant Green
    C_EMERALD_LIGHT = RGBColor(236, 253, 245)  # #ECFDF5 Green Tint
    C_AMBER         = RGBColor(245, 158, 11)   # #F59E0B Warm Amber
    C_AMBER_LIGHT   = RGBColor(254, 243, 199)  # #FEF3C7 Amber Tint
    C_PURPLE        = RGBColor(139, 92, 246)   # #8B5CF6 Deep Violet
    C_PURPLE_LIGHT  = RGBColor(245, 243, 255)  # #F5F3FF Violet Tint
    C_ROSE          = RGBColor(239, 68, 68)    # #EF4444 Alert Rose
    C_ROSE_LIGHT    = RGBColor(254, 242, 242)  # #FEF2F2 Rose Tint

    # Text Colors
    C_TEXT_DARK     = RGBColor(15, 23, 42)     # #0F172A Bold Dark
    C_TEXT_MUTED    = RGBColor(71, 85, 105)    # #475569 Subtitle Slate
    C_TEXT_WHITE    = RGBColor(255, 255, 255)  # #FFFFFF Pure White
    C_TEXT_LIGHT_MUTED = RGBColor(148, 163, 184) # #94A3B8 Dimmed Light

    FONT = "Segoe UI"

    # ── Helpers ─────────────────────────────────────────────────────────────
    def set_bg(slide, color):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_header(slide, badge, title, subtitle=None, is_dark=False):
        # Category Badge
        badge_w = Inches(max(2.6, len(badge) * 0.13))
        badge_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(0.48), badge_w, Inches(0.34))
        badge_box.fill.solid()
        badge_box.fill.fore_color.rgb = RGBColor(30, 58, 138) if is_dark else C_PRIMARY_LIGHT
        badge_box.line.color.rgb = C_CYAN if is_dark else C_PRIMARY
        badge_box.line.width = Pt(1)
        tf_b = badge_box.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
        p_b = tf_b.paragraphs[0]
        p_b.alignment = PP_ALIGN.CENTER
        p_b.text = badge.upper()
        p_b.font.name = FONT
        p_b.font.size = Pt(10)
        p_b.font.bold = True
        p_b.font.color.rgb = C_CYAN if is_dark else C_PRIMARY

        # Title & Subtitle Box
        tb = slide.shapes.add_textbox(Inches(0.9), Inches(0.92), Inches(11.533), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = C_TEXT_WHITE if is_dark else C_TEXT_DARK

        if subtitle:
            p_sub = tf.add_paragraph()
            p_sub.text = subtitle
            p_sub.font.name = FONT
            p_sub.font.size = Pt(13.5)
            p_sub.font.color.rgb = C_TEXT_LIGHT_MUTED if is_dark else C_TEXT_MUTED
            p_sub.space_before = Pt(3)

    def add_footer(slide, slide_num, total_slides=13, is_dark=False):
        tb = slide.shapes.add_textbox(Inches(0.9), Inches(6.92), Inches(8.5), Inches(0.35))
        tf = tb.text_frame
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = "ResearchBridge • Final Presentation • Institute of Information Technology (IIT), University of Dhaka"
        p.font.name = FONT
        p.font.size = Pt(9.5)
        p.font.color.rgb = C_TEXT_LIGHT_MUTED if is_dark else C_TEXT_MUTED

        tb_num = slide.shapes.add_textbox(Inches(10.433), Inches(6.92), Inches(2.0), Inches(0.35))
        tf_num = tb_num.text_frame
        tf_num.margin_left = tf_num.margin_top = tf_num.margin_right = tf_num.margin_bottom = 0
        p_num = tf_num.paragraphs[0]
        p_num.alignment = PP_ALIGN.RIGHT
        p_num.text = f"{slide_num:02d} / {total_slides:02d}"
        p_num.font.name = FONT
        p_num.font.size = Pt(10.5)
        p_num.font.bold = True
        p_num.font.color.rgb = C_CYAN if is_dark else C_PRIMARY

    def add_card(slide, left, top, width, height, title=None, badge=None, accent_color=C_PRIMARY, is_dark=False):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = C_DARK_CARD if is_dark else C_LIGHT_CARD
        card.line.color.rgb = C_DARK_BORDER if is_dark else C_LIGHT_BORDER
        card.line.width = Pt(1.5)

        # Top Accent Strip
        strip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.08))
        strip.fill.solid()
        strip.fill.fore_color.rgb = accent_color
        strip.line.fill.background()

        tb = slide.shapes.add_textbox(left + Inches(0.24), top + Inches(0.2), width - Inches(0.48), height - Inches(0.3))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        if badge:
            p_badge = tf.paragraphs[0]
            p_badge.text = badge.upper()
            p_badge.font.name = FONT
            p_badge.font.size = Pt(9.5)
            p_badge.font.bold = True
            p_badge.font.color.rgb = accent_color
            p_title = tf.add_paragraph()
        else:
            p_title = tf.paragraphs[0]

        if title:
            p_title.text = title
            p_title.font.name = FONT
            p_title.font.size = Pt(15)
            p_title.font.bold = True
            p_title.font.color.rgb = C_TEXT_WHITE if is_dark else C_TEXT_DARK
            if badge:
                p_title.space_before = Pt(2)

        return tf

    def add_notes(slide, timing_str, script_str):
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = f"⏱️ TIME: {timing_str}\n\n🎙️ WHAT TO SAY:\n{script_str}"

    card_w3 = Inches(3.68)
    gap3    = Inches(0.24)
    start_x3= Inches(0.9)
    card_top= Inches(1.95)
    card_h  = Inches(4.68)

    card_w4 = Inches(2.72)
    gap4    = Inches(0.2)
    start_x4= Inches(0.9)

    # =========================================================================
    # SLIDE 1: TITLE SLIDE (Executive Dark Presentation Cover)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_bg(slide1, C_DARK_BG)

    # Institution Badge Banner
    badge_cover = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(0.8), Inches(6.8), Inches(0.42))
    badge_cover.fill.solid()
    badge_cover.fill.fore_color.rgb = RGBColor(30, 58, 138)
    badge_cover.line.color.rgb = C_CYAN
    badge_cover.line.width = Pt(1)
    tf_bc = badge_cover.text_frame
    tf_bc.margin_left = tf_bc.margin_top = tf_bc.margin_right = tf_bc.margin_bottom = 0
    p_bc = tf_bc.paragraphs[0]
    p_bc.alignment = PP_ALIGN.CENTER
    p_bc.text = "INSTITUTE OF INFORMATION TECHNOLOGY • UNIVERSITY OF DHAKA"
    p_bc.font.name = FONT
    p_bc.font.size = Pt(11)
    p_bc.font.bold = True
    p_bc.font.color.rgb = C_CYAN

    # Title Box
    t_box = slide1.shapes.add_textbox(Inches(0.9), Inches(1.4), Inches(11.533), Inches(2.6))
    tf1 = t_box.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_top = tf1.margin_right = tf1.margin_bottom = 0

    p_main = tf1.paragraphs[0]
    p_main.text = "ResearchBridge"
    p_main.font.name = FONT
    p_main.font.size = Pt(46)
    p_main.font.bold = True
    p_main.font.color.rgb = C_TEXT_WHITE
    p_main.space_after = Pt(6)

    p_sub = tf1.add_paragraph()
    p_sub.text = "The AI-Powered Research Launchpad for Emerging Scholars"
    p_sub.font.name = FONT
    p_sub.font.size = Pt(20)
    p_sub.font.bold = True
    p_sub.font.color.rgb = C_CYAN
    p_sub.space_after = Pt(6)

    p_tag = tf1.add_paragraph()
    p_tag.text = "Connecting Minds • Grounding Scientific Discovery • Accelerating Publication"
    p_tag.font.name = FONT
    p_tag.font.size = Pt(14)
    p_tag.font.color.rgb = C_TEXT_LIGHT_MUTED

    # Presenter Card
    card_p = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(4.35), Inches(5.55), Inches(2.2))
    card_p.fill.solid()
    card_p.fill.fore_color.rgb = C_DARK_CARD
    card_p.line.color.rgb = C_DARK_BORDER
    card_p.line.width = Pt(1.5)
    tf_p = card_p.text_frame
    tf_p.word_wrap = True
    tf_p.margin_left = tf_p.margin_top = tf_p.margin_right = tf_p.margin_bottom = Inches(0.24)
    
    p = tf_p.paragraphs[0]
    p.text = "PRESENTER"
    p.font.name = FONT
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = C_CYAN
    
    p2 = tf_p.add_paragraph()
    p2.text = "Mostofa Aminur Rashid"
    p2.font.name = FONT
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = C_TEXT_WHITE
    p2.space_before = Pt(3)

    p3 = tf_p.add_paragraph()
    p3.text = "Roll: 2506107  •  IIT, University of Dhaka"
    p3.font.name = FONT
    p3.font.size = Pt(13)
    p3.font.color.rgb = C_TEXT_LIGHT_MUTED
    p3.space_before = Pt(2)

    # Supervisor Card
    card_s = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.88), Inches(4.35), Inches(5.55), Inches(2.2))
    card_s.fill.solid()
    card_s.fill.fore_color.rgb = C_DARK_CARD
    card_s.line.color.rgb = C_DARK_BORDER
    card_s.line.width = Pt(1.5)
    tf_s = card_s.text_frame
    tf_s.word_wrap = True
    tf_s.margin_left = tf_s.margin_top = tf_s.margin_right = tf_s.margin_bottom = Inches(0.24)
    
    p = tf_s.paragraphs[0]
    p.text = "PROJECT SUPERVISOR"
    p.font.name = FONT
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = C_EMERALD
    
    p2 = tf_s.add_paragraph()
    p2.text = "Dr. Kazi Muheymin-Us-Sakib"
    p2.font.name = FONT
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = C_TEXT_WHITE
    p2.space_before = Pt(3)

    p3 = tf_s.add_paragraph()
    p3.text = "Professor  •  IIT, University of Dhaka"
    p3.font.name = FONT
    p3.font.size = Pt(13)
    p3.font.color.rgb = C_TEXT_LIGHT_MUTED
    p3.space_before = Pt(2)

    add_notes(slide1, "0:00 - 0:15 (15 Seconds)",
              "Respected teachers and evaluators, Good day! I am Mostofa Aminur Rashid, presenting 'ResearchBridge — The AI-Powered Research Launchpad for Emerging Scholars', supervised by Dr. Kazi Muheymin-Us-Sakib, Professor at IIT, University of Dhaka.")

    # =========================================================================
    # SLIDE 2: WHAT IS THE PROJECT? (Clear 3-Pillar Unified Platform)
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_bg(slide2, C_LIGHT_BG)
    add_header(slide2, "1. Executive Overview", "What is ResearchBridge?", "A single, unified platform that guides researchers from early concept to published paper.")
    add_footer(slide2, 2)

    pillars = [
        ("Smart Discovery",
         "Find the right collaborators, mentors, and foundational research papers by understanding the deeper meaning of your ideas, not just matching keywords.",
         "✔ Solves the 0-citation cold start\n✔ Replaces cold emails & manual search\n✔ Connects junior & senior researchers",
         C_PRIMARY, "01 • DISCOVER"),
        ("Live Workspace",
         "Write and edit research papers together in real-time with zero typing conflicts, integrated sprint task boards, and a shared team PDF library.",
         "✔ Real-time multi-author co-authoring\n✔ Centralized team milestone tracking\n✔ Integrated PDF notes & highlights",
         C_CYAN, "02 • EXECUTE"),
        ("Publication AI",
         "Conduct thorough literature reviews with an AI assistant that provides 100% verified citations, matches top target journals, and audits formatting.",
         "✔ 100% cited answers from real papers\n✔ AI-powered journal recommendation\n✔ Reduces manuscript desk rejections",
         C_EMERALD, "03 • PUBLISH")
    ]

    for i, (title, desc, bullets, color, badge) in enumerate(pillars):
        x = start_x3 + i * (card_w3 + gap3)
        tf_c = add_card(slide2, x, card_top, card_w3, card_h, title, badge=badge, accent_color=color)
        
        p_desc = tf_c.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = FONT
        p_desc.font.size = Pt(13.5)
        p_desc.font.color.rgb = C_TEXT_MUTED
        p_desc.space_before = Pt(8)

        p_b = tf_c.add_paragraph()
        p_b.text = f"\n{bullets}"
        p_b.font.name = FONT
        p_b.font.size = Pt(12.5)
        p_b.font.bold = True
        p_b.font.color.rgb = C_TEXT_DARK
        p_b.space_before = Pt(12)

    add_notes(slide2, "0:15 - 0:40 (25 Seconds)",
              "What is ResearchBridge? Simply put: it replaces 5 fragmented tools with one end-to-end platform for academic research. It provides (1) Smart Discovery to find the right partners and papers, (2) A Live Workspace for seamless team writing, and (3) Publication AI that generates 100% cited literature reviews and matches target journals.")

    # =========================================================================
    # SLIDE 3: THE PROBLEM (Clear Real-World Academic Bottlenecks)
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_bg(slide3, C_LIGHT_BG)
    add_header(slide3, "2. Problem Statement", "The 3 Core Bottlenecks in Research Today", "Why emerging scholars and student researchers struggle to get published.")
    add_footer(slide3, 3)

    problems = [
        ("Academic Isolation",
         "Emerging scholars start with 0 citations and no established network. Existing sites like ResearchGate focus almost exclusively on famous professors, leaving newcomers stranded.",
         "❌ No way for beginners to get discovered\n❌ Painful manual networking & cold emails\n❌ Good research ideas die in isolation",
         "BOTTLENECK 1", C_ROSE),
        ("Tool Fragmentation",
         "Research teams are forced to juggle 5+ disjointed tools: Slack for chat, Google Docs or Overleaf for writing, Zotero for PDFs, and Trello for tasks. Context constantly gets lost.",
         "❌ Version conflicts & overwritten drafts\n❌ Scattered files and lost citations\n❌ Massive administrative overhead",
         "BOTTLENECK 2", C_AMBER),
        ("Publishing Frustration",
         "Over 60% of papers get desk-rejected because of wrong journal scope. Meanwhile, general AI tools like ChatGPT frequently fabricate fake citations that ruin academic credibility.",
         "❌ High rejection rates from poor journal fit\n❌ AI tools hallucinate fake sources\n❌ Weeks wasted on manual re-formatting",
         "BOTTLENECK 3", C_PURPLE)
    ]

    for i, (title, desc, bullets, badge, color) in enumerate(problems):
        x = start_x3 + i * (card_w3 + gap3)
        tf_c = add_card(slide3, x, card_top, card_w3, card_h, title, badge=badge, accent_color=color)
        
        p_desc = tf_c.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = FONT
        p_desc.font.size = Pt(13.5)
        p_desc.font.color.rgb = C_TEXT_MUTED
        p_desc.space_before = Pt(8)

        p_b = tf_c.add_paragraph()
        p_b.text = f"\n{bullets}"
        p_b.font.name = FONT
        p_b.font.size = Pt(12.5)
        p_b.font.bold = True
        p_b.font.color.rgb = color
        p_b.space_before = Pt(12)

    add_notes(slide3, "0:40 - 1:05 (25 Seconds)",
              "Why is this needed? Today's researchers face three major bottlenecks: First, Isolation — beginners have zero citation track record and no network. Second, Tool Fragmentation — jumping between Docs, Overleaf, Slack, and Zotero wastes hours and creates version conflicts. Third, Publishing Frustration — high desk rejection rates and general AI that makes up fake citations.")

    # =========================================================================
    # SLIDE 4: HOW IT WORKS (Intuitive 4-Step Human Lifecycle)
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_bg(slide4, C_LIGHT_BG)
    add_header(slide4, "3. User Journey", "How ResearchBridge Works in 4 Steps", "A clear and guided journey from initial research interest to final journal submission.")
    add_footer(slide4, 4)

    steps = [
        ("Step 1: Input Topic",
         "Describe your research topic or questions in plain language.",
         "⚡ AI immediately matches compatible co-researchers, faculty mentors, and foundation literature.",
         "🤝 INSTANT MATCH", C_PRIMARY, "PHASE 01"),
        ("Step 2: Launch Workspace",
         "Create a shared project room with one click.",
         "⚡ Co-author drafts together in real-time with zero typing conflicts and track milestone tasks.",
         "✍️ LIVE CO-AUTHORING", C_CYAN, "PHASE 02"),
        ("Step 3: Grounded AI Review",
         "Explore related literature using the built-in AI research assistant.",
         "⚡ Get accurate synthesized answers where every single sentence cites a real paper in the database.",
         "📚 100% VERIFIED CITATIONS", C_EMERALD, "PHASE 03"),
        ("Step 4: Match & Publish",
         "Prepare your manuscript for successful submission.",
         "⚡ Auto-generate instant citations (BibTeX/APA), match indexed journals, and audit submission compliance.",
         "🚀 SUBMISSION READY", C_PURPLE, "PHASE 04")
    ]

    for i, (title, desc, outcome, badge_text, color, phase_badge) in enumerate(steps):
        x = start_x4 + i * (card_w4 + gap4)
        tf_c = add_card(slide4, x, card_top, card_w4, card_h, title, badge=phase_badge, accent_color=color)
        
        p_desc = tf_c.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = FONT
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = C_TEXT_MUTED
        p_desc.space_before = Pt(8)

        p_out = tf_c.add_paragraph()
        p_out.text = f"\n{outcome}"
        p_out.font.name = FONT
        p_out.font.size = Pt(12)
        p_out.font.color.rgb = C_TEXT_DARK
        p_out.space_before = Pt(8)

        p_badge = tf_c.add_paragraph()
        p_badge.text = f"\n{badge_text}"
        p_badge.font.name = FONT
        p_badge.font.size = Pt(12.5)
        p_badge.font.bold = True
        p_badge.font.color.rgb = color
        p_badge.space_before = Pt(14)

    add_notes(slide4, "1:05 - 1:25 (20 Seconds)",
              "How does a researcher use the platform? Step 1: Input your research interests to match partners and papers. Step 2: Open a live collaborative workspace. Step 3: Use grounded AI for trustworthy literature synthesis. Step 4: Find the right target journal and run pre-submission compliance checks.")

    # =========================================================================
    # SLIDE 5: SYSTEM ARCHITECTURE (Intuitive 3-Tier Layered Design)
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_bg(slide5, C_LIGHT_BG)
    add_header(slide5, "4. System Architecture", "Modern 3-Tier Microservices Architecture", "Designed for high responsiveness, real-time collaboration, and continuous reliability.")
    add_footer(slide5, 5)

    arch_tiers = [
        ("Interactive Frontend",
         "Modern User Interface\n(Next.js & React)\n\n• Ultra-fast, responsive web interface\n• Real-time live cursor co-editing\n• Dynamic Kanban milestone tracking\n• Seamless mobile & desktop design",
         C_PRIMARY, "TIER 1 • USER EXPERIENCE"),
        ("Secure API Gateway",
         "Central API & Security Hub\n(Node.js & Express)\n\n• Enterprise JWT authentication & security\n• Real-time WebSocket connection router\n• Intelligent rate limiting & data shield\n• Fast event distribution between services",
         C_CYAN, "TIER 2 • API & SECURITY"),
        ("AI Intelligence Core",
         "Machine Learning Service\n(Python & FastAPI)\n\n• Deep semantic understanding (SBERT)\n• Grounded literature question-answering\n• Automatic fallback engine (zero crashes)\n• Smart journal scope matching",
         C_EMERALD, "TIER 3 • AI ENGINE")
    ]

    for i, (title, body, color, badge) in enumerate(arch_tiers):
        x = start_x3 + i * (card_w3 + gap3)
        tf_c = add_card(slide5, x, Inches(1.95), card_w3, Inches(3.25), title, badge=badge, accent_color=color)
        p_b = tf_c.add_paragraph()
        p_b.text = body
        p_b.font.name = FONT
        p_b.font.size = Pt(12.5)
        p_b.font.color.rgb = C_TEXT_DARK
        p_b.space_before = Pt(6)

    # Polyglot Database Banner
    db_card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(5.35), Inches(11.533), Inches(1.35))
    db_card.fill.solid()
    db_card.fill.fore_color.rgb = C_DARK_BG
    db_card.line.color.rgb = C_DARK_BORDER
    db_card.line.width = Pt(1.5)
    tf_db = db_card.text_frame
    tf_db.word_wrap = True
    tf_db.margin_left = tf_db.margin_top = tf_db.margin_right = tf_db.margin_bottom = Inches(0.18)
    
    p_dh = tf_db.paragraphs[0]
    p_dh.text = "INTEGRATED SPECIALIZED DATA LAYER (PURPOSE-BUILT PERSISTENCE)"
    p_dh.font.name = FONT
    p_dh.font.size = Pt(10.5)
    p_dh.font.bold = True
    p_dh.font.color.rgb = C_CYAN

    p_dc = tf_db.add_paragraph()
    p_dc.text = "🐘 PostgreSQL (Core Data)   •   ⚡ Elasticsearch (AI Semantic Search)   •   🕸️ Neo4j (Trust Network)   •   🚀 Redis (High-Speed Cache)   •   📦 MinIO (PDF Library)"
    p_dc.font.name = FONT
    p_dc.font.size = Pt(13)
    p_dc.font.bold = True
    p_dc.font.color.rgb = C_TEXT_WHITE
    p_dc.space_before = Pt(4)

    add_notes(slide5, "1:25 - 1:45 (20 Seconds)",
              "Under the hood, the system is built in 3 modular layers: An ultra-responsive Next.js frontend, a secure Node.js API Gateway, and a dedicated Python AI service for semantic matching and RAG, powered by purpose-built databases for relational data, vector search, trust graphs, and caching.")

    # =========================================================================
    # SLIDE 6: CORE FEATURE 1 — GROUNDED AI ASSISTANT (No Hallucinations)
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_bg(slide6, C_DARK_BG)
    add_header(slide6, "Core Innovation 1", "Grounded AI Assistant with 100% Real Citations", "Scholarly AI generation that strictly references verified literature and never hallucinates.", is_dark=True)
    add_footer(slide6, 6, is_dark=True)

    w1_cards = [
        ("100% Real Citations",
         "Unlike ChatGPT which frequently invents non-existent papers, our assistant strictly grounds every claim in verified literature, providing clickable inline references [1], [2].",
         "✔ Eliminates fake citations\n✔ Clickable links to real papers\n✔ Maintains complete academic integrity",
         C_EMERALD, "INTEGRITY & TRUST"),
        ("Focused Reasoning",
         "Uses a controlled 2-round retrieval loop that understands user intent, expands relevant keywords, and synthesizes answers in under two seconds without getting lost.",
         "✔ Sub-2-second response time\n✔ Understands contextual intent\n✔ Clear synthesized summaries",
         C_CYAN, "SPEED & ACCURACY"),
        ("Failsafe Zero-Crash Engine",
         "If external cloud AI providers experience outages or network hiccups, the system automatically switches to instant extractive summarization so researchers are never blocked.",
         "✔ 100% platform availability\n✔ Zero server crash errors\n✔ Works reliably even offline",
         C_AMBER, "RELIABILITY SHIELD")
    ]

    for i, (title, desc, bullets, color, badge) in enumerate(w1_cards):
        x = start_x3 + i * (card_w3 + gap3)
        tf_c = add_card(slide6, x, card_top, card_w3, card_h, title, badge=badge, accent_color=color, is_dark=True)
        
        p_desc = tf_c.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = FONT
        p_desc.font.size = Pt(13.5)
        p_desc.font.color.rgb = C_TEXT_LIGHT_MUTED
        p_desc.space_before = Pt(8)

        p_b = tf_c.add_paragraph()
        p_b.text = f"\n{bullets}"
        p_b.font.name = FONT
        p_b.font.size = Pt(12.5)
        p_b.font.bold = True
        p_b.font.color.rgb = C_TEXT_WHITE
        p_b.space_before = Pt(12)

    add_notes(slide6, "1:45 - 2:10 (25 Seconds)",
              "Our flagship innovation is our Grounded AI Assistant. Unlike generic chatbots that invent fake references, our engine provides 100% verified citations directly linked to real papers. Plus, our built-in failsafe fallback guarantees zero server crashes even if external cloud AI goes down.")

    # =========================================================================
    # SLIDE 7: CORE FEATURE 2 — TRUSTRANK (Fair Academic Reputation)
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_bg(slide7, C_LIGHT_BG)
    add_header(slide7, "Core Innovation 2", "TrustRank: Merit-Based Academic Credibility", "Creating a fair launchpad for young researchers beyond superficial citation counts.")
    add_footer(slide7, 7)

    w2_cards = [
        ("Holistic Credibility",
         "Traditional academic metrics only count citations, ignoring promising beginners. TrustRank measures active mentorship, peer validations, and verified community contributions.",
         "✔ Breaks the zero-citation barrier\n✔ Rewards active research efforts\n✔ Values mentorship & guidance",
         C_PURPLE, "FAIR REPUTATION"),
        ("Trust Flow Network",
         "Credibility starts with verified professors and naturally propagates to active student mentees as they collaborate, solve problems, and contribute quality work.",
         "✔ Faculty endorsement power\n✔ Verified junior researchers rise\n✔ Real-time reputation scoring",
         C_PRIMARY, "MERITOCRATIC FLOW"),
        ("Anti-Spam Protection",
         "Intelligent graph analysis detects and blocks fake accounts, artificial upvote rings, and spam, ensuring academic discussions remain constructive and trustworthy.",
         "✔ Immunity to fake upvote rings\n✔ Sub-15ms graph verification\n✔ Connects trusted 2nd-degree peers",
         C_EMERALD, "COMMUNITY INTEGRITY")
    ]

    for i, (title, desc, bullets, color, badge) in enumerate(w2_cards):
        x = start_x3 + i * (card_w3 + gap3)
        tf_c = add_card(slide7, x, card_top, card_w3, card_h, title, badge=badge, accent_color=color)
        
        p_desc = tf_c.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = FONT
        p_desc.font.size = Pt(13.5)
        p_desc.font.color.rgb = C_TEXT_MUTED
        p_desc.space_before = Pt(8)

        p_b = tf_c.add_paragraph()
        p_b.text = f"\n{bullets}"
        p_b.font.name = FONT
        p_b.font.size = Pt(12.5)
        p_b.font.bold = True
        p_b.font.color.rgb = C_TEXT_DARK
        p_b.space_before = Pt(12)

    add_notes(slide7, "2:10 - 2:30 (20 Seconds)",
              "Our second key innovation is TrustRank. Traditional metrics penalize beginners who haven't published yet. Our TrustRank propagates credibility from verified professors through real mentorship and community contributions, creating a true meritocracy.")

    # =========================================================================
    # SLIDE 8: CORE FEATURE 3 — REAL-TIME COLLABORATIVE WORKSPACE
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_bg(slide8, C_LIGHT_BG)
    add_header(slide8, "Core Innovation 3", "Real-Time Collaborative Workspace", "Seamless co-authoring, agile project task tracking, and a shared team paper library.")
    add_footer(slide8, 8)

    w3_cards = [
        ("Conflict-Free Co-Authoring",
         "Multiple co-authors can type, edit, and format papers simultaneously in real-time. The system automatically reconciles edits with zero merge conflicts or lost work.",
         "✔ Sub-100ms real-time typing sync\n✔ Multi-cursor live presence\n✔ Seamless offline auto-reconnect",
         C_PRIMARY, "LIVE WRITING"),
        ("Agile Milestone Boards",
         "Visual Kanban sprint boards allow teams to organize paper sections, attach reference PDFs directly to tasks, and hit submission deadlines smoothly.",
         "✔ Replaces separate Trello/Notion\n✔ Links tasks directly to drafts\n✔ Clear accountability & progress",
         C_CYAN, "TEAM MANAGEMENT"),
        ("Shared Team Paper Library",
         "A unified team PDF repository where members can highlight text, add shared notes, extract key figures, and build a collective knowledge base.",
         "✔ Instant PDF text & metadata extraction\n✔ Shared team highlights & notes\n✔ One-click citation export",
         C_EMERALD, "KNOWLEDGE BASE")
    ]

    for i, (title, desc, bullets, color, badge) in enumerate(w3_cards):
        x = start_x3 + i * (card_w3 + gap3)
        tf_c = add_card(slide8, x, card_top, card_w3, card_h, title, badge=badge, accent_color=color)
        
        p_desc = tf_c.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = FONT
        p_desc.font.size = Pt(13.5)
        p_desc.font.color.rgb = C_TEXT_MUTED
        p_desc.space_before = Pt(8)

        p_b = tf_c.add_paragraph()
        p_b.text = f"\n{bullets}"
        p_b.font.name = FONT
        p_b.font.size = Pt(12.5)
        p_b.font.bold = True
        p_b.font.color.rgb = C_TEXT_DARK
        p_b.space_before = Pt(12)

    add_notes(slide8, "2:30 - 2:50 (20 Seconds)",
              "Our third capability is the Real-Time Collaborative Workspace. Co-authors can draft papers simultaneously with sub-100ms sync and zero conflicts, unified with Kanban milestone tracking and a shared team paper library.")

    # =========================================================================
    # SLIDE 9: CORE FEATURE 4 — PUBLICATION ASSISTANT & JOURNAL MATCHMAKER
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_bg(slide9, C_LIGHT_BG)
    add_header(slide9, "Core Innovation 4", "AI Journal Matchmaker & Publication Suite", "Guiding research papers from first draft to successful journal acceptance.")
    add_footer(slide9, 9)

    w4_cards = [
        ("Smart Journal Matcher",
         "Scans your draft abstract and matches it against thousands of verified DOAJ and Scimago indexed journals to find the ideal scope, impact factor, and turnaround speed.",
         "✔ Matches verified Q1-Q4 journals\n✔ Filters Open Access & review speeds\n✔ Maximizes chances of acceptance",
         C_AMBER, "JOURNAL DISCOVERY"),
        ("1-Click Citation Suite",
         "Instantly formats references into BibTeX, APA, IEEE, and ACM standards. Automatically pulls citation data from DOIs and uploaded PDFs without manual typing.",
         "✔ Zero manual reference formatting\n✔ Instant BibTeX / APA / IEEE export\n✔ Seamless sync to workspace editor",
         C_PRIMARY, "CITATION AUTOMATION"),
        ("Pre-Submission Compliance",
         "Audits manuscript structure (IMRAD), ethics statements, methodology details, and formatting standards before you submit to minimize desk rejections.",
         "✔ Catches missing structural sections\n✔ Audits citation completeness\n✔ Prevents preventable rejections",
         C_PURPLE, "QUALITY ASSURANCE")
    ]

    for i, (title, desc, bullets, color, badge) in enumerate(w4_cards):
        x = start_x3 + i * (card_w3 + gap3)
        tf_c = add_card(slide9, x, card_top, card_w3, card_h, title, badge=badge, accent_color=color)
        
        p_desc = tf_c.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = FONT
        p_desc.font.size = Pt(13.5)
        p_desc.font.color.rgb = C_TEXT_MUTED
        p_desc.space_before = Pt(8)

        p_b = tf_c.add_paragraph()
        p_b.text = f"\n{bullets}"
        p_b.font.name = FONT
        p_b.font.size = Pt(12.5)
        p_b.font.bold = True
        p_b.font.color.rgb = C_TEXT_DARK
        p_b.space_before = Pt(12)

    add_notes(slide9, "2:50 - 3:10 (20 Seconds)",
              "Our fourth core capability is the Publication Assistant. When your paper is drafted, our AI matcher recommends the right DOAJ and Scimago journals, formats instant citations, and runs pre-submission checks to prevent desk rejections.")

    # =========================================================================
    # SLIDE 10: COMPETITIVE ADVANTAGE (Why ResearchBridge Wins)
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_bg(slide10, C_LIGHT_BG)
    add_header(slide10, "5. Platform Comparison", "Why ResearchBridge Outperforms Existing Tools", "A side-by-side comparison showing how ResearchBridge closes critical academic gaps.")
    add_footer(slide10, 10)

    comps = [
        ("ResearchGate",
         "Static Profile Directory",
         "❌ Caters only to senior professors\n❌ No live co-authoring editor\n❌ No beginner matchmaking\n❌ No citation AI assistant",
         RGBColor(148, 163, 184)),
        ("Google Scholar",
         "Passive Search Engine",
         "❌ Passive search index only\n❌ Zero collaboration or team tools\n❌ No mentor connections\n❌ No journal matchmaker",
         RGBColor(148, 163, 184)),
        ("Overleaf / Slack",
         "Generic Fragmented Tools",
         "❌ No academic discovery AI\n❌ No verified citation grounding\n❌ No journal recommendation\n❌ High subscription costs",
         RGBColor(148, 163, 184)),
        ("ResearchBridge",
         "Unified AI Research Ecosystem",
         "✔ Smart semantic partner matching\n✔ Sub-100ms live co-authoring\n✔ 100% cited literature AI\n✔ Built-in journal matchmaker",
         C_PRIMARY)
    ]

    for i, (p_title, p_sub, p_points, color) in enumerate(comps):
        x = start_x4 + i * (card_w4 + gap4)
        is_highlight = (i == 3)
        tf_c = add_card(slide10, x, card_top, card_w4, card_h, p_title, badge="EXISTING TOOL" if not is_highlight else "OUR PLATFORM", accent_color=color, is_dark=is_highlight)
        
        p_s = tf_c.add_paragraph()
        p_s.text = p_sub
        p_s.font.name = FONT
        p_s.font.size = Pt(12)
        p_s.font.color.rgb = C_TEXT_LIGHT_MUTED if is_highlight else C_TEXT_MUTED
        p_s.space_before = Pt(4)

        p_pts = tf_c.add_paragraph()
        p_pts.text = f"\n{p_points}"
        p_pts.font.name = FONT
        p_pts.font.size = Pt(12.5)
        p_pts.font.bold = True
        p_pts.font.color.rgb = C_CYAN if is_highlight else (C_ROSE if "❌" in p_points else C_TEXT_DARK)
        p_pts.space_before = Pt(10)

    add_notes(slide10, "3:10 - 3:25 (15 Seconds)",
              "Compared to existing tools: ResearchGate only serves established authors; Google Scholar is just a passive search index; Overleaf and Slack lack academic intelligence. ResearchBridge is the ONLY unified ecosystem that connects discovery, live writing, and publication.")

    # =========================================================================
    # SLIDE 11: BENCHMARK RESULTS (Bold Clear Performance Metrics)
    # =========================================================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_bg(slide11, C_LIGHT_BG)
    add_header(slide11, "6. Validation & Performance", "Empirical Benchmarks & System Responsiveness", "Rigorous performance testing validates smooth execution under high concurrent load.")
    add_footer(slide11, 11)

    stats = [
        ("< 42 ms", "AI Search Speed",
         "High-dimensional semantic search retrieves the most relevant papers and researchers in milliseconds.",
         "Industry Benchmark: < 100ms", C_PRIMARY, "AI LATENCY"),
        ("> 1,450", "Requests / Second",
         "API Gateway sustains heavy concurrent load without bottlenecking or dropping requests.",
         "Target Capacity: > 1,000 req/s", C_EMERALD, "THROUGHPUT"),
        ("< 100 ms", "Live Sync Latency",
         "Real-time co-authoring document updates synchronize instantaneously between active team members.",
         "Acceptable UX Limit: < 200ms", C_CYAN, "REAL-TIME SYNC"),
        ("100 %", "Query Availability",
         "Deterministic failsafe fallback ensures zero server crashes and reliable answers even if cloud AI goes down.",
         "High-Availability Target: 99.9%", C_PURPLE, "RELIABILITY")
    ]

    for i, (val, title, desc, target, color, badge) in enumerate(stats):
        x = start_x4 + i * (card_w4 + gap4)
        tf_c = add_card(slide11, x, card_top, card_w4, card_h, title, badge=badge, accent_color=color)
        
        # Big Stat
        p_val = tf_c.add_paragraph()
        p_val.text = val
        p_val.font.name = FONT
        p_val.font.size = Pt(30)
        p_val.font.bold = True
        p_val.font.color.rgb = color
        p_val.space_before = Pt(4)

        p_desc = tf_c.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = FONT
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = C_TEXT_MUTED
        p_desc.space_before = Pt(6)

        p_tgt = tf_c.add_paragraph()
        p_tgt.text = f"\n✅ {target}"
        p_tgt.font.name = FONT
        p_tgt.font.size = Pt(11.5)
        p_tgt.font.bold = True
        p_tgt.font.color.rgb = C_EMERALD
        p_tgt.space_before = Pt(10)

    add_notes(slide11, "3:25 - 3:45 (20 Seconds)",
              "We validated the platform under heavy load: AI vector search completes in under 42ms; our API gateway sustains over 1,450 requests per second; real-time live typing sync is below 100ms; and our RAG fallback achieves 100% query reliability with zero crashes.")

    # =========================================================================
    # SLIDE 12: FUTURE ROADMAP (Strategic Horizons)
    # =========================================================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_bg(slide12, C_LIGHT_BG)
    add_header(slide12, "7. Future Scope", "Future Horizons for ResearchBridge", "Strategic expansion into inter-university networks, multimodal AI, and autonomous review.")
    add_footer(slide12, 12)

    roadmap = [
        ("Campus Federation",
         "Connecting Universities\n\n• Integration with eduroam & institutional SSO\n• Verified cross-university researcher badges\n• Inter-department research collaboration hubs\n• Automated campus onboarding",
         C_PRIMARY, "HORIZON 1 • EXPANSION"),
        ("Multimodal Lab AI",
         "Voice & Visual Intelligence\n\n• Real-time voice agent for lab meeting minutes\n• Automatic diagram, table & formula extraction\n• Integrated video conference workrooms\n• Automated code notebook synchronization",
         C_CYAN, "HORIZON 2 • MULTIMODAL"),
        ("Autonomous Review",
         "Pre-Submission Critique AI\n\n• Adversarial pre-submission paper review\n• Pinpoints methodology and logic gaps\n• Automated journal-specific LaTeX compilation\n• Citation freshness and retraction checks",
         C_PURPLE, "HORIZON 3 • AUTONOMOUS")
    ]

    for i, (title, body, color, badge) in enumerate(roadmap):
        x = start_x3 + i * (card_w3 + gap3)
        tf_c = add_card(slide12, x, card_top, card_w3, card_h, title, badge=badge, accent_color=color)
        p_b = tf_c.add_paragraph()
        p_b.text = body
        p_b.font.name = FONT
        p_b.font.size = Pt(13.5)
        p_b.font.color.rgb = C_TEXT_DARK
        p_b.space_before = Pt(8)

    add_notes(slide12, "3:45 - 3:55 (10 Seconds)",
              "Looking ahead, our roadmap focuses on 3 horizons: Inter-university campus federation with single-sign-on, multimodal voice and visual lab assistants, and an autonomous AI peer-review engine for pre-submission critiques.")

    # =========================================================================
    # SLIDE 13: CONCLUSION & Q&A (High Impact Executive Close)
    # =========================================================================
    slide13 = prs.slides.add_slide(blank_layout)
    set_bg(slide13, C_DARK_BG)
    add_header(slide13, "Conclusion & Summary", "Empowering the Next Generation of Scholars", "Bridging minds, grounding scientific discovery, and accelerating scholarly innovation.", is_dark=True)
    add_footer(slide13, 13, is_dark=True)

    summary_cards = [
        ("Unified Ecosystem",
         "Replaces 5+ fragmented apps with a single, intuitive launchpad covering discovery, execution, and publication.",
         C_CYAN, "01 • UNIFIED WORKFLOW"),
        ("Trustworthy AI",
         "Grounded AI with 100% real literature citations and zero server crashes through intelligent failsafe fallback.",
         C_EMERALD, "02 • CITATION RIGOR"),
        ("Meritocratic Growth",
         "TrustRank breaks the 0-citation cold start, empowering emerging researchers to gain reputation and publish faster.",
         C_AMBER, "03 • FAIR REPUTATION")
    ]

    for i, (title, body, color, badge) in enumerate(summary_cards):
        x = start_x3 + i * (card_w3 + gap3)
        tf_c = add_card(slide13, x, Inches(1.95), card_w3, Inches(3.1), title, badge=badge, accent_color=color, is_dark=True)
        p_b = tf_c.add_paragraph()
        p_b.text = body
        p_b.font.name = FONT
        p_b.font.size = Pt(13.5)
        p_b.font.color.rgb = C_TEXT_LIGHT_MUTED
        p_b.space_before = Pt(8)

    # Thank you & Contact Box
    end_box = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(5.35), Inches(11.533), Inches(1.35))
    end_box.fill.solid()
    end_box.fill.fore_color.rgb = C_DARK_CARD
    end_box.line.color.rgb = C_CYAN
    end_box.line.width = Pt(1.5)
    tf_end = end_box.text_frame
    tf_end.word_wrap = True
    tf_end.margin_left = tf_end.margin_top = tf_end.margin_right = tf_end.margin_bottom = Inches(0.15)

    p_e1 = tf_end.paragraphs[0]
    p_e1.alignment = PP_ALIGN.CENTER
    p_e1.text = "THANK YOU! • QUESTIONS & ANSWERS"
    p_e1.font.name = FONT
    p_e1.font.size = Pt(16)
    p_e1.font.bold = True
    p_e1.font.color.rgb = C_CYAN

    p_e2 = tf_end.add_paragraph()
    p_e2.alignment = PP_ALIGN.CENTER
    p_e2.text = "Presenter: Mostofa Aminur Rashid (Roll: 2506107)   |   Supervisor: Dr. Kazi Muheymin-Us-Sakib\nInstitute of Information Technology (IIT), University of Dhaka"
    p_e2.font.name = FONT
    p_e2.font.size = Pt(12.5)
    p_e2.font.color.rgb = C_TEXT_LIGHT_MUTED
    p_e2.space_before = Pt(4)

    add_notes(slide13, "3:55 - 4:10 (15 Seconds)",
              "In conclusion, ResearchBridge turns isolated individual efforts into collaborative academic breakthroughs. Thank you very much for your time and attention. I am now ready for your questions!")

    # Save presentation
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    out = os.path.abspath(r"c:\project\github\SmartResearch\SmartResearch_Final_Presentation.pptx")
    create_deck(out)
